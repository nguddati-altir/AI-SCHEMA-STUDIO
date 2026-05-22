"""Generate SchemaDesigner.pptx — pitch deck for the Schema Designer hackathon project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Palette — matches the app's dark UI
BG = RGBColor(0x0F, 0x17, 0x2A)        # slate-900
PANEL = RGBColor(0x11, 0x18, 0x27)     # slate-800
BORDER = RGBColor(0x33, 0x41, 0x55)
TEXT = RGBColor(0xE5, 0xE7, 0xEB)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)    # sky
ACCENT2 = RGBColor(0x81, 0x8C, 0xF8)   # indigo
OK = RGBColor(0x10, 0xB9, 0x81)
WARN = RGBColor(0xF5, 0x9E, 0x0B)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tb


def add_panel(slide, x, y, w, h, *, fill=PANEL, line=BORDER, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    return shape


def add_header(slide, title, subtitle=None):
    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
        title, size=32, bold=True, color=ACCENT,
    )
    if subtitle:
        add_text(
            slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.5),
            subtitle, size=16, color=MUTED,
        )
    # accent rule
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.55), Inches(1.2), Pt(2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT2
    rule.line.fill.background()


def add_footer(slide, page, total):
    add_text(
        slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
        "Schema Designer  ·  Hackathon 2026", size=10, color=MUTED,
    )
    add_text(
        slide, Inches(12.2), Inches(7.05), Inches(1.0), Inches(0.3),
        f"{page} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def bullet_list(slide, x, y, w, h, items, *, size=18, color=TEXT, leading=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(leading)
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color


def add_code(slide, x, y, w, h, code, *, size=12):
    panel = add_panel(slide, x, y, w, h, fill=RGBColor(0x0B, 0x12, 0x20), line=BORDER, radius=True)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = TEXT


TOTAL_SLIDES = 11


# ---------- Slide 1 — Title ----------
s = add_slide()
# big gradient-like band
band = slide_band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0), SW, Inches(3.5))
band.fill.solid()
band.fill.fore_color.rgb = PANEL
band.line.fill.background()

add_text(s, Inches(0.8), Inches(2.4), Inches(12), Inches(1.4),
         "Schema Designer", size=64, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(3.6), Inches(12), Inches(0.7),
         "Visual schema designer for MySQL / Postgres — with AI assist and Excel seeding.",
         size=22, color=TEXT)
add_text(s, Inches(0.85), Inches(4.3), Inches(12), Inches(0.5),
         "React + FastAPI  ·  Prisma + SQL output  ·  Claude Opus 4.7",
         size=16, color=ACCENT2)
add_text(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
         "Hackathon 2026", size=11, color=MUTED)


# ---------- Slide 2 — Problem ----------
s = add_slide()
add_header(s, "The problem", "Starting a new project shouldn't mean writing DDL by hand.")
bullet_list(s, Inches(0.6), Inches(2.0), Inches(12), Inches(4),
            [
                "Every new project starts with the same boilerplate: tables, columns, FKs, JSON fields, seed data.",
                "Hand-writing CREATE TABLE for two dialects (MySQL + Postgres) is repetitive and error-prone.",
                "Prisma schema and SQL DDL drift apart over time — the source of truth becomes unclear.",
                "Seeding tables from a customer's Excel file means a one-off import script every single time.",
                "And nobody wants to type out 20 columns for a table they could describe in one sentence.",
            ], size=18, leading=12)
add_footer(s, 2, TOTAL_SLIDES)


# ---------- Slide 3 — What it does ----------
s = add_slide()
add_header(s, "What it does", "Three things, one tool.")

box_y = Inches(2.0)
box_h = Inches(4.5)
box_w = Inches(4.0)
gap = Inches(0.25)

cards = [
    ("Design",
     "Visual table designer.\n\nDrag-free editor for tables, columns, types (incl. JSON), PK / FK / unique / auto-increment / defaults.\n\nSchema auto-persists to localStorage.",
     ACCENT),
    ("Generate",
     "One source → four outputs.\n\nPostgres DDL · MySQL DDL · Prisma schema (PG) · Prisma schema (MySQL).\n\nLive preview, copy or download.",
     ACCENT2),
    ("Seed",
     "Excel → INSERT statements.\n\nUpload .xlsx (one sheet per table, headers must match). Validates, previews, then emits dialect-specific INSERTs.",
     OK),
]
x = Inches(0.6)
for title, body, color in cards:
    panel = add_panel(s, x, box_y, box_w, box_h, radius=True)
    add_text(s, x + Inches(0.3), box_y + Inches(0.3), box_w - Inches(0.6), Inches(0.6),
             title, size=24, bold=True, color=color)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.3), box_y + Inches(0.95), Inches(0.7), Pt(2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()
    add_text(s, x + Inches(0.3), box_y + Inches(1.2), box_w - Inches(0.6), box_h - Inches(1.5),
             body, size=15, color=TEXT)
    x += box_w + gap

add_footer(s, 3, TOTAL_SLIDES)


# ---------- Slide 4 — AI Assist (the special sauce) ----------
s = add_slide()
add_header(s, "✨ AI Assist", "Describe an entity in plain English. Get suggested columns.")

# Left: feature list
bullet_list(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5),
            [
                "Per-table AI panel — \"A user account with profile JSON and timestamps\"",
                "Returns a list of suggested columns with type, attributes, and reasoning",
                "Check / uncheck each suggestion — you stay in control",
                "Follow-up prompts: \"add soft-delete\", \"make email lowercase-unique\"",
                "\"New table from description\" — generate a whole table from one sentence",
                "Uses Claude Opus 4.7 with adaptive thinking + prompt caching",
            ], size=16, leading=10)

# Right: mock UI panel
ax, ay, aw, ah = Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.5)
add_panel(s, ax, ay, aw, ah, radius=True)
add_text(s, ax + Inches(0.25), ay + Inches(0.2), aw - Inches(0.5), Inches(0.4),
         "✨ AI Assist", size=14, bold=True, color=ACCENT)

prompt_box = add_panel(s, ax + Inches(0.25), ay + Inches(0.75), aw - Inches(0.5), Inches(0.65),
                       fill=RGBColor(0x0B, 0x12, 0x20), radius=True)
add_text(s, ax + Inches(0.4), ay + Inches(0.85), aw - Inches(0.8), Inches(0.5),
         "A user account with email, password hash, profile JSON, and timestamps",
         size=11, color=MUTED, font="Consolas")

# Suggested columns
items = [
    ("id", "int", "PK · AI · NN", ACCENT),
    ("email", "string(255)", "UNQ · NN", ACCENT2),
    ("password_hash", "string(255)", "NN", ACCENT2),
    ("profile", "json", "—", OK),
    ("created_at", "datetime", "= CURRENT_TIMESTAMP", ACCENT2),
    ("updated_at", "datetime", "—", ACCENT2),
]
ry = ay + Inches(1.55)
for name, ty, attrs, col in items:
    add_text(s, ax + Inches(0.35), ry, Inches(1.6), Inches(0.3), f"✓  {name}", size=11, color=ACCENT, font="Consolas")
    add_text(s, ax + Inches(1.95), ry, Inches(1.3), Inches(0.3), ty, size=11, color=TEXT, font="Consolas")
    add_text(s, ax + Inches(3.25), ry, Inches(2.4), Inches(0.3), attrs, size=10, color=MUTED, font="Consolas")
    ry += Inches(0.35)

add_text(s, ax + Inches(0.35), ay + ah - Inches(0.45), aw - Inches(0.7), Inches(0.3),
         "6 of 6 selected   →   [ Add 6 columns ]", size=11, color=ACCENT2)

add_footer(s, 4, TOTAL_SLIDES)


# ---------- Slide 5 — Architecture ----------
s = add_slide()
add_header(s, "Architecture", "Two services, clean boundary.")

# Frontend box
fx, fy, fw, fh = Inches(0.8), Inches(2.2), Inches(5.4), Inches(4.0)
add_panel(s, fx, fy, fw, fh, radius=True)
add_text(s, fx + Inches(0.3), fy + Inches(0.25), fw - Inches(0.6), Inches(0.5),
         "Frontend  ·  React + Vite + TS", size=18, bold=True, color=ACCENT)
bullet_list(s, fx + Inches(0.4), fy + Inches(0.9), fw - Inches(0.6), fh - Inches(1.0),
            [
                "App.tsx — 3 views: Design / Generate / Import",
                "TableEditor — per-table column editor + FK picker",
                "AIAssist — describe → suggest → accept/refine",
                "GeneratedCode — live SQL & Prisma preview",
                "ExcelImport — validate + INSERT generator",
                "localStorage persistence, Vite proxy → backend",
            ], size=13, leading=6)

# Backend box
bx, by, bw, bh = Inches(7.1), Inches(2.2), Inches(5.4), Inches(4.0)
add_panel(s, bx, by, bw, bh, radius=True)
add_text(s, bx + Inches(0.3), by + Inches(0.25), bw - Inches(0.6), Inches(0.5),
         "Backend  ·  FastAPI (Python)", size=18, bold=True, color=ACCENT2)
bullet_list(s, bx + Inches(0.4), by + Inches(0.9), bw - Inches(0.6), bh - Inches(1.0),
            [
                "models.py — Pydantic Schema / Table / Column",
                "generators/sql.py — Postgres + MySQL DDL",
                "generators/prisma.py — schema.prisma w/ relations",
                "excel.py — openpyxl validator + INSERT generator",
                "ai.py — Claude Opus 4.7 via official SDK",
                "main.py — REST endpoints + CORS",
            ], size=13, leading=6)

# Arrow between them
arrow = s.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, Inches(6.25), Inches(4.05), Inches(0.85), Inches(0.35))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT
arrow.line.fill.background()
add_text(s, Inches(6.0), Inches(4.4), Inches(1.4), Inches(0.3),
         "/api → :8000", size=10, color=MUTED, align=PP_ALIGN.CENTER, font="Consolas")

add_footer(s, 5, TOTAL_SLIDES)


# ---------- Slide 6 — Type mapping ----------
s = add_slide()
add_header(s, "One schema, every flavor", "11 column types map cleanly to Postgres, MySQL, and Prisma.")

headers = ["Schema type", "Postgres", "MySQL", "Prisma"]
rows = [
    ("int",       "INTEGER / SERIAL",  "INT [AUTO_INCREMENT]", "Int"),
    ("bigint",    "BIGINT / BIGSERIAL","BIGINT",               "BigInt"),
    ("string(n)", "VARCHAR(n)",        "VARCHAR(n)",           "String @db.VarChar(n)"),
    ("text",      "TEXT",              "TEXT",                 "String @db.Text"),
    ("boolean",   "BOOLEAN",           "TINYINT(1)",           "Boolean"),
    ("decimal",   "NUMERIC(p,s)",      "DECIMAL(p,s)",         "Decimal"),
    ("datetime",  "TIMESTAMP",         "DATETIME",             "DateTime"),
    ("json",      "JSONB",             "JSON",                 "Json"),
    ("uuid",      "UUID",              "CHAR(36)",             "String @db.Uuid"),
]

tx, ty = Inches(0.6), Inches(2.0)
col_w = [Inches(2.0), Inches(3.2), Inches(3.4), Inches(3.6)]
row_h = Inches(0.42)

# Header row
hx = tx
for i, ht in enumerate(headers):
    panel = add_panel(s, hx, ty, col_w[i], row_h, fill=PANEL, line=BORDER)
    add_text(s, hx + Inches(0.15), ty + Inches(0.08), col_w[i] - Inches(0.3), row_h,
             ht, size=12, bold=True, color=ACCENT)
    hx += col_w[i]

# Body rows
ry = ty + row_h
for row in rows:
    hx = tx
    for i, cell in enumerate(row):
        add_panel(s, hx, ry, col_w[i], row_h, fill=BG, line=BORDER)
        add_text(s, hx + Inches(0.15), ry + Inches(0.08), col_w[i] - Inches(0.3), row_h,
                 cell, size=11, color=TEXT, font="Consolas")
        hx += col_w[i]
    ry += row_h

add_footer(s, 6, TOTAL_SLIDES)


# ---------- Slide 7 — Generated output sample ----------
s = add_slide()
add_header(s, "Generated output", "Pick a tab. Copy or download.")

# Left — Postgres DDL
sql_code = """-- Schema: my_project
-- Dialect: postgres

CREATE TABLE "users" (
  "id"         SERIAL NOT NULL,
  "email"      VARCHAR(255) NOT NULL UNIQUE,
  "name"       VARCHAR(120),
  "metadata"   JSONB,
  "created_at" TIMESTAMP NOT NULL DEFAULT CURRENT_TIMESTAMP,
  PRIMARY KEY ("id")
);

CREATE TABLE "posts" (
  "id"      SERIAL NOT NULL,
  "user_id" INTEGER NOT NULL,
  "title"   VARCHAR(200) NOT NULL,
  "tags"    JSONB,
  PRIMARY KEY ("id"),
  FOREIGN KEY ("user_id") REFERENCES "users"("id")
);"""

prisma_code = """model users {
  id         Int       @id @default(autoincrement())
  email      String    @unique @db.VarChar(255)
  name       String?   @db.VarChar(120)
  metadata   Json?
  created_at DateTime  @db.Timestamp(6)
  posts_list posts[]
}

model posts {
  id           Int    @id @default(autoincrement())
  user_id      Int
  title        String @db.VarChar(200)
  tags         Json?
  users_user_id users  @relation(fields: [user_id], references: [id])
}"""

add_text(s, Inches(0.6), Inches(1.95), Inches(6), Inches(0.4),
         "Postgres SQL", size=14, bold=True, color=ACCENT)
add_code(s, Inches(0.6), Inches(2.4), Inches(6.1), Inches(4.4), sql_code, size=11)

add_text(s, Inches(6.95), Inches(1.95), Inches(6), Inches(0.4),
         "Prisma (Postgres)", size=14, bold=True, color=ACCENT2)
add_code(s, Inches(6.95), Inches(2.4), Inches(5.75), Inches(4.4), prisma_code, size=11)

add_footer(s, 7, TOTAL_SLIDES)


# ---------- Slide 8 — Excel import flow ----------
s = add_slide()
add_header(s, "Excel → seed data", "Headers must match. The rest is automatic.")

# Flow diagram
fy = Inches(2.3)
fh = Inches(1.3)

stages = [
    ("📤  Upload", ".xlsx workbook\nOne sheet per table", ACCENT),
    ("✓  Validate", "Sheet name = table\nHeaders = column names", ACCENT2),
    ("👁  Preview", "First 25 rows per sheet\nFlag missing/extra cols", OK),
    ("⚙  Generate", "INSERT statements\nfor Postgres or MySQL", WARN),
]
fw = Inches(2.85)
gap = Inches(0.2)
fx = Inches(0.6)
for i, (head, body, color) in enumerate(stages):
    panel = add_panel(s, fx, fy, fw, fh, radius=True)
    panel.line.color.rgb = color
    panel.line.width = Pt(1.5)
    add_text(s, fx + Inches(0.2), fy + Inches(0.15), fw - Inches(0.4), Inches(0.5),
             head, size=15, bold=True, color=color)
    add_text(s, fx + Inches(0.2), fy + Inches(0.65), fw - Inches(0.4), Inches(0.7),
             body, size=11, color=TEXT)
    fx += fw + gap
    if i < len(stages) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   fx - gap - Inches(0.05), fy + Inches(0.55),
                                   gap + Inches(0.05), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MUTED
        arrow.line.fill.background()

# Sample INSERTs
sample = """-- Generated INSERT statements (postgres)

-- users
INSERT INTO "users" ("email", "name", "metadata")
  VALUES ('alice@example.com', 'Alice', '{"role":"admin"}');
INSERT INTO "users" ("email", "name", "metadata")
  VALUES ('bob@example.com',   'Bob',   '{"role":"user"}');

-- posts
INSERT INTO "posts" ("user_id", "title", "tags")
  VALUES (1, 'Hello world', '["intro","welcome"]');"""

add_text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.4),
         "Output", size=12, bold=True, color=MUTED)
add_code(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.4), sample, size=11)

add_footer(s, 8, TOTAL_SLIDES)


# ---------- Slide 9 — AI deep-dive ----------
s = add_slide()
add_header(s, "Under the hood — AI Assist", "How the magic actually works.")

# Diagram: prompt → Claude → structured output → UI
boxes = [
    ("User", "Describes entity\nor refines prior turn", ACCENT),
    ("Backend", "Loads system prompt\n+ conversation history\n+ existing columns", ACCENT2),
    ("Claude Opus 4.7", "Adaptive thinking\nPrompt-cached system\nclient.messages.parse()", OK),
    ("Pydantic", "Schema-validated\nSuggestedColumn[]\n+ rationale", WARN),
    ("UI", "Checkboxes per col\n+ accept/refine loop", ACCENT),
]
bw = Inches(2.4)
gap = Inches(0.05)
by = Inches(2.2)
bh = Inches(1.6)
bx = Inches(0.6)
for i, (h, b, color) in enumerate(boxes):
    p = add_panel(s, bx, by, bw, bh, radius=True)
    p.line.color.rgb = color
    p.line.width = Pt(1.5)
    add_text(s, bx + Inches(0.15), by + Inches(0.15), bw - Inches(0.3), Inches(0.5),
             h, size=14, bold=True, color=color)
    add_text(s, bx + Inches(0.15), by + Inches(0.6), bw - Inches(0.3), Inches(1.0),
             b, size=10, color=TEXT)
    if i < len(boxes) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   bx + bw, by + Inches(0.65), Inches(0.15), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MUTED
        arrow.line.fill.background()
    bx += bw + Inches(0.15)

# Code snippet
ai_code = """response = client.messages.parse(
    model="claude-opus-4-7",
    max_tokens=8000,
    thinking={"type": "adaptive"},
    system=[{"type": "text", "text": SYSTEM_PROMPT,
             "cache_control": {"type": "ephemeral"}}],
    messages=[*history, {"role": "user", "content": user_turn}],
    output_format=SuggestionResult,   # Pydantic — schema auto-derived
)
result = response.parsed_output       # typed list of SuggestedColumn"""

add_text(s, Inches(0.6), Inches(4.1), Inches(12), Inches(0.4),
         "The Claude API call", size=12, bold=True, color=MUTED)
add_code(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.3), ai_code, size=11)

add_footer(s, 9, TOTAL_SLIDES)


# ---------- Slide 10 — Run it ----------
s = add_slide()
add_header(s, "Run it on Windows", "Two terminals. Five minutes.")

backend_cmd = """# Terminal 1 — backend
cd backend
python -m venv .venv
.\\.venv\\Scripts\\Activate.ps1
pip install -r requirements.txt

# Optional — enable AI Assist
$env:ANTHROPIC_API_KEY = "sk-ant-..."

uvicorn app.main:app --reload --port 8000"""

frontend_cmd = """# Terminal 2 — frontend
cd frontend
npm install
npm run dev

# Open in browser
# http://127.0.0.1:5173"""

add_text(s, Inches(0.6), Inches(1.95), Inches(6), Inches(0.4),
         "Backend  ·  FastAPI", size=14, bold=True, color=ACCENT)
add_code(s, Inches(0.6), Inches(2.4), Inches(6.1), Inches(4.0), backend_cmd, size=13)

add_text(s, Inches(6.95), Inches(1.95), Inches(6), Inches(0.4),
         "Frontend  ·  Vite", size=14, bold=True, color=ACCENT2)
add_code(s, Inches(6.95), Inches(2.4), Inches(5.75), Inches(4.0), frontend_cmd, size=13)

add_text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.4),
         "API docs at http://127.0.0.1:8000/docs   ·   Schema persists to localStorage   ·   Works without API key (AI panel just disables)",
         size=11, color=MUTED)

add_footer(s, 10, TOTAL_SLIDES)


# ---------- Slide 11 — What's next ----------
s = add_slide()
add_header(s, "What's next", "Hackathon scope shipped. Here's the roadmap.")

next_items = [
    ("Live DB connect",   "Run generated DDL directly against a Postgres/MySQL connection string from the UI."),
    ("Indexes & checks",  "Add CREATE INDEX, CHECK constraints, and composite-unique indexes to the designer."),
    ("Schema import",     "Reverse direction: introspect an existing database and pre-populate the designer."),
    ("More ORMs",         "TypeORM, SQLAlchemy, Drizzle output alongside Prisma."),
    ("Diff & migrations", "Compare two schemas and emit ALTER TABLE migrations."),
    ("Team mode",         "Share schemas via a backend store; comment threads on columns."),
]

cy = Inches(2.0)
ch = Inches(0.75)
for title, body in next_items:
    add_text(s, Inches(0.6), cy, Inches(3.0), ch,
             title, size=16, bold=True, color=ACCENT2)
    add_text(s, Inches(3.7), cy, Inches(9.0), ch,
             body, size=14, color=TEXT)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), cy + Inches(0.65), Inches(12), Pt(0.5))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BORDER
    rule.line.fill.background()
    cy += ch

# Closer
add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.4),
         "Thanks!  ·  github.com/your-org/schema-designer", size=14, bold=True,
         color=ACCENT, align=PP_ALIGN.CENTER)

add_footer(s, 11, TOTAL_SLIDES)


out = "SchemaDesigner.pptx"
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
